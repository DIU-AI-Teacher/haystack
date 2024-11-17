from fastapi import FastAPI, File, UploadFile, HTTPException, Query
from fastapi.middleware.cors import CORSMiddleware
from haystack.document_stores import ElasticsearchDocumentStore
from haystack.nodes import PreProcessor, FARMReader, ElasticsearchRetriever
from haystack.pipelines import ExtractiveQAPipeline
from haystack.utils import convert_files_to_docs
import uvicorn
import os
import shutil
from typing import List, Optional
from pydantic import BaseModel
import PyPDF2
from pptx import Presentation
import json

app = FastAPI(title="AI Professor API")

# Enable CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize document store
document_store = ElasticsearchDocumentStore(
    host="localhost",
    port=9200,
    username="",
    password="",
    index="course_content"
)

# Initialize components
preprocessor = PreProcessor(
    clean_empty_lines=True,
    clean_whitespace=True,
    clean_header_footer=True,
    split_by="word",
    split_length=500,
    split_overlap=50
)

# Initialize Reader & Retriever
reader = FARMReader(model_name_or_path="deepset/roberta-base-squad2")
retriever = ElasticsearchRetriever(document_store=document_store)

# Create pipeline
pipe = ExtractiveQAPipeline(reader, retriever)

# Create storage directories
os.makedirs("uploads", exist_ok=True)
os.makedirs("processed", exist_ok=True)

class CourseContent(BaseModel):
    title: str
    content: str
    file_type: str
    useful_links: Optional[List[str]] = []

@app.post("/upload-content/")
async def upload_content(
    file: UploadFile = File(...),
    course_title: str = Query(..., description="Title of the course"),
    useful_links: str = Query("[]", description="JSON array of useful links")
):
    try:
        # Save uploaded file
        file_path = f"uploads/{file.filename}"
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # Extract text based on file type
        text_content = ""
        file_type = file.filename.split('.')[-1].lower()

        if file_type == 'pdf':
            with open(file_path, 'rb') as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                for page in pdf_reader.pages:
                    text_content += page.extract_text()
        
        elif file_type == 'pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"

        # Preprocess and index the content
        docs = preprocessor.process([{
            "content": text_content,
            "meta": {
                "course_title": course_title,
                "file_name": file.filename,
                "file_type": file_type,
                "useful_links": json.loads(useful_links)
            }
        }])
        
        document_store.write_documents(docs)

        return {"message": "Content uploaded and indexed successfully"}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/search/")
async def search(
    question: str = Query(..., description="Question to ask about the course content"),
    course_title: Optional[str] = Query(None, description="Filter by course title")
):
    try:
        # Prepare filters
        filters = {}
        if course_title:
            filters["course_title"] = course_title

        # Get prediction
        prediction = pipe.run(
            query=question,
            params={
                "Retriever": {"top_k": 3, "filters": filters},
                "Reader": {"top_k": 1}
            }
        )

        # Extract answer and context
        if prediction['answers']:
            answer = prediction['answers'][0]
            
            # Get useful links from the document metadata
            useful_links = []
            if prediction['documents']:
                useful_links = prediction['documents'][0].meta.get('useful_links', [])

            return {
                "answer": answer.answer,
                "context": answer.context,
                "confidence": answer.score,
                "useful_links": useful_links
            }
        else:
            return {
                "answer": "I couldn't find an answer to your question in the course materials.",
                "context": None,
                "confidence": 0.0,
                "useful_links": []
            }

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/courses/")
async def list_courses():
    try:
        # Get all documents from the document store
        documents = document_store.get_all_documents()
        
        # Extract unique course titles
        courses = set()
        for doc in documents:
            courses.add(doc.meta.get('course_title'))
        
        return {"courses": list(courses)}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
